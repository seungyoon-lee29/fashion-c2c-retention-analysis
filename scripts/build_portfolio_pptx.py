"""Build a portfolio PPTX for the fashion C2C retention project."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT_DIR = ROOT / "portfolio"
OUT = OUT_DIR / "Lee_Seungyoon_Product_DA_Portfolio_Project1.pptx"
FIG = ROOT / "docs" / "figures"
DATA_SOURCE = "Source: Mercari MerRec (`mercari-us/merrec`), 20230501 parquet x5, local cache `data/raw`"

INK = RGBColor(28, 32, 38)
MUTED = RGBColor(91, 102, 118)
GREEN = RGBColor(31, 79, 70)
TEAL = RGBColor(15, 118, 110)
NAVY = RGBColor(28, 43, 67)
AMBER = RGBColor(180, 83, 9)
RED = RGBColor(185, 28, 28)
BG = RGBColor(248, 249, 250)
LINE = RGBColor(220, 226, 232)
WHITE = RGBColor(255, 255, 255)


def add_textbox(slide, x, y, w, h, text, size=18, color=INK, bold=False, align=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Inches(0.04)
    tf.margin_right = Inches(0.04)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.text = text
    if align:
        p.alignment = align
    run = p.runs[0]
    run.font.name = "Aptos"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_title(slide, title, subtitle=None):
    add_textbox(slide, 0.55, 0.32, 8.7, 0.38, title, 22, INK, True)
    if subtitle:
        add_textbox(slide, 0.58, 0.75, 8.3, 0.28, subtitle, 9.5, MUTED)
    add_textbox(slide, 11.35, 0.38, 1.35, 0.22, "Project 01", 8.5, MUTED, False, PP_ALIGN.RIGHT)
    line = slide.shapes.add_shape(1, Inches(0.55), Inches(1.05), Inches(12.2), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = LINE
    line.line.fill.background()


def add_footer(slide):
    add_textbox(slide, 0.55, 7.02, 4.0, 0.2, "이승윤 | Product / Growth Data Analyst", 7.5, MUTED)
    add_textbox(slide, 8.35, 7.02, 4.4, 0.2, "github.com/seungyoon-lee29/fashion-c2c-retention-analysis", 7.5, MUTED, False, PP_ALIGN.RIGHT)


def add_card(slide, x, y, w, h, title, body, accent=TEAL, value=None):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = LINE
    shape.line.width = Pt(0.8)
    bar = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(0.07), Inches(h))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()
    if value:
        add_textbox(slide, x + 0.18, y + 0.15, w - 0.35, 0.36, value, 21, accent, True)
        add_textbox(slide, x + 0.18, y + 0.57, w - 0.35, 0.22, title, 9.5, INK, True)
        add_textbox(slide, x + 0.18, y + 0.86, w - 0.35, h - 0.96, body, 8.5, MUTED)
    else:
        add_textbox(slide, x + 0.18, y + 0.16, w - 0.35, 0.22, title, 11, INK, True)
        add_textbox(slide, x + 0.18, y + 0.48, w - 0.35, h - 0.55, body, 8.8, MUTED)
    return shape


def add_bullets(slide, x, y, w, h, bullets, size=12, color=INK, gap=0.12):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.space_after = Pt(gap * 12)
        run = p.runs[0]
        run.font.name = "Aptos"
        run.font.size = Pt(size)
        run.font.color.rgb = color
    return box


def add_section_label(slide, x, y, label, color=GREEN):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(1.22), Inches(0.28))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    add_textbox(slide, x + 0.08, y + 0.065, 1.05, 0.12, label, 7.5, WHITE, True, PP_ALIGN.CENTER)


def add_table_like(slide, x, y, w, rows, col_widths, header=True):
    row_h = 0.42
    for r, row in enumerate(rows):
        yy = y + r * row_h
        xx = x
        fill = RGBColor(238, 247, 245) if r == 0 and header else WHITE
        for c, cell in enumerate(row):
            cw = w * col_widths[c]
            shape = slide.shapes.add_shape(1, Inches(xx), Inches(yy), Inches(cw), Inches(row_h))
            shape.fill.solid()
            shape.fill.fore_color.rgb = fill
            shape.line.color.rgb = LINE
            add_textbox(slide, xx + 0.05, yy + 0.09, cw - 0.1, row_h - 0.1, str(cell), 8.0, INK if r == 0 else MUTED, r == 0)
            xx += cw


def blank(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BG
    return slide


def build():
    OUT_DIR.mkdir(exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 1
    s = blank(prs)
    shape = s.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = GREEN
    shape.line.fill.background()
    add_section_label(s, 0.75, 0.74, "PORTFOLIO", NAVY)
    add_textbox(s, 0.78, 1.35, 7.8, 0.45, "패션 C2C 신규 유저 리텐션 분석", 30, WHITE, True)
    add_textbox(s, 0.82, 1.95, 7.8, 0.35, "관측 로그를 A/B 실험 설계로 연결한 Product/Growth DA 대표 프로젝트", 14, RGBColor(224, 242, 238))
    add_textbox(s, 0.82, 5.85, 5.7, 0.28, "이승윤 | Product / Growth Data Analyst", 12, WHITE, True)
    add_textbox(s, 0.82, 6.22, 6.7, 0.22, "sy3097@gmail.com  |  github.com/seungyoon-lee29", 9.5, RGBColor(224, 242, 238))
    add_textbox(s, 0.82, 6.55, 7.4, 0.2, DATA_SOURCE, 7.2, RGBColor(224, 242, 238))
    add_card(
        s,
        8.5,
        1.35,
        3.7,
        1.05,
        "Dataset",
        "Mercari MerRec C2C behavior logs, 20230501 parquet x5",
        TEAL,
        "2.77M events",
    )
    add_card(s, 8.5, 2.62, 3.7, 1.05, "Users", "New observed user journey", TEAL, "43,311 users")
    add_card(s, 8.5, 3.89, 3.7, 1.05, "Final handoff", "First-session exploration breadth nudge A/B", AMBER, "A/B design")

    # 2
    s = blank(prs)
    add_title(s, "Executive Summary", "구매가 아니라 D7 재방문을 실험 KPI로 전환")
    add_card(s, 0.65, 1.35, 2.85, 1.45, "구매완료율", "구매 전환만으로는 신규 유저 경험 의사결정이 어렵다.", AMBER, "4.4%")
    add_card(s, 3.75, 1.35, 2.85, 1.45, "D7 재방문", "실험의 primary KPI로 전환했다.", TEAL, "50.1%")
    add_card(s, 6.85, 1.35, 2.85, 1.45, "코호트 gap", "첫날 활동 폭 상위-하위 관측 차이. 인과효과 아님.", TEAL, "+27pp")
    add_card(s, 9.95, 1.35, 2.65, 1.45, "실험 규모", "MDE 2.0pp 기준 필요 표본.", NAVY, "9,805/arm")
    add_textbox(s, 0.75, 3.35, 4.2, 0.32, "핵심 판단", 16, INK, True)
    add_bullets(
        s,
        0.8,
        3.85,
        5.15,
        2.1,
        [
            "구매완료율 4.4%는 희소해 신규 유저 경험의 주지표로 약하다.",
            "D7 재방문 50.1%를 primary KPI로 두면 더 안정적으로 실험할 수 있다.",
            "첫날 활동 폭은 강한 후보 신호지만, 관측 데이터만으로 인과효과를 확정하지 않는다.",
        ],
        12,
    )
    add_card(
        s,
        6.55,
        3.35,
        5.65,
        2.45,
        "결론",
        "카트 최적화가 아니라 첫 세션에서 인접 카테고리와 저장 후보를 노출하는 탐색 폭 넛지를 A/B로 검증한다.",
        GREEN,
        "First-session exploration breadth nudge",
    )
    add_footer(s)

    # 3
    s = blank(prs)
    add_title(s, "Data QA Gate", "분석 가능한 로그인지 먼저 검증")
    add_card(s, 0.7, 1.4, 2.7, 1.15, "Status", "Phase-0 data quality gate", TEAL, "PASS")
    add_card(s, 3.65, 1.4, 2.7, 1.15, "Time window", "2023-05-01 ~ 2023-05-31", NAVY, "30 days")
    add_card(s, 6.6, 1.4, 2.7, 1.15, "Events", "Expected vocabulary only", TEAL, "6 types")
    add_card(s, 9.55, 1.4, 2.7, 1.15, "Price QA", "No negative or zero prices", NAVY, "median 25.0")
    add_textbox(s, 0.76, 2.7, 11.5, 0.22, DATA_SOURCE, 7.5, MUTED)
    rows = [
        ("Check", "Result", "Portfolio signal"),
        ("Data source", "mercari-us/merrec 20230501 parquet x5", "출처 추적 가능"),
        ("Required fields", "user_id, stime, session_id, event null 없음", "분석 단위 신뢰"),
        ("Event vocabulary", "6개 expected events 외 값 없음", "퍼널 정의 가능"),
        ("Timestamp", "2023-05-01 ~ 2023-05-31", "관찰 윈도우 명확"),
        ("Known warning", "반복 로그 키 2,466건", "session_id는 user-scoped로 해석"),
    ]
    add_table_like(s, 0.72, 3.05, 11.85, rows, [0.22, 0.44, 0.34])
    add_footer(s)

    # 4
    s = blank(prs)
    add_title(s, "Funnel Diagnosis", "최대 유출은 조회에서 찜으로 넘어가는 구간")
    steps = [
        ("view", "98.7%"),
        ("like", "52.3%"),
        ("cart", "22.9%"),
        ("offer", "17.1%"),
        ("buy_start", "5.7%"),
        ("buy_comp", "4.4%"),
    ]
    x0 = 0.75
    for i, (name, pct) in enumerate(steps):
        x = x0 + i * 2.02
        color = AMBER if name == "buy_comp" else TEAL if i < 2 else NAVY
        add_card(s, x, 1.55, 1.55, 1.05, name, "user reach", color, pct)
        if i < len(steps) - 1:
            add_textbox(s, x + 1.55, 1.93, 0.35, 0.2, "→", 16, MUTED, True, PP_ALIGN.CENTER)
    add_card(s, 0.75, 3.25, 5.45, 1.7, "Largest leak", "조회 -> 찜 구간에서 직전 도달 후 미도달 20,082명. 조회 직후 저장 후보 또는 인접 카테고리 노출이 실험 후보로 연결된다.", AMBER, "20,082 users")
    add_card(s, 6.55, 3.25, 5.65, 1.7, "Decision", "구매완료율 4.4%는 너무 희소하다. 신규 유저 경험의 북극성은 구매 전환이 아니라 D7 재방문으로 둔다.", GREEN, "KPI switch")
    add_footer(s)

    # 5
    s = blank(prs)
    add_title(s, "Cohort Signal", "첫날 경험 폭이 넓은 유저가 더 돌아온다")
    add_card(s, 0.75, 1.35, 3.3, 1.4, "Baseline", "신규 관측 유저 D7 재방문", TEAL, "50.1%")
    add_card(s, 4.35, 1.35, 3.3, 1.4, "Q1 first-day activity", "첫날 활동량 하위 코호트 D7", NAVY, "39.3%")
    add_card(s, 7.95, 1.35, 3.3, 1.4, "Q4 first-day activity", "첫날 활동량 상위 코호트 D7", GREEN, "66.2%")
    add_card(s, 0.75, 3.35, 4.9, 1.65, "Observed gap", "첫날 활동 폭 상위-하위 코호트의 D7 재방문 차이. 이 값은 목표 uplift가 아니라 후보 신호다.", AMBER, "+27pp")
    add_card(s, 6.0, 3.35, 5.25, 1.65, "Interpretation", "첫날 경험 폭은 실험 후보로 충분히 강하다. 하지만 유저의 원래 적극성이 섞여 있으므로 인과효과는 A/B에서만 판단한다.", GREEN, "Signal, not proof")
    add_footer(s)

    # 6
    s = blank(prs)
    add_title(s, "Identification Check", "좋아 보이는 레버를 그대로 믿지 않았다")
    rows = [
        ("Candidate", "Naive association", "Adjusted diagnostic", "Decision"),
        ("cart>=1", "+16.0pp", "+3.2pp", "양성 확정 불가"),
        ("cat>=3", "+25.4pp", "-6.5pp", "보정 후 불안정"),
        ("like>=3 / offer>=1", "강한 신호", "overlap 부족", "점추정 철회"),
    ]
    add_table_like(s, 0.8, 1.45, 11.2, rows, [0.28, 0.22, 0.24, 0.26])
    add_card(s, 0.85, 3.75, 5.25, 1.55, "What changed after adjustment", "겹침이 충분한 레버도 보정 후 0 근처 또는 음수로 약해졌다. 관측 신호를 바로 제품 레버로 해석하면 과대주장 위험이 있다.", RED, "Causal honesty")
    if (FIG / "impact_estimators.png").exists():
        s.shapes.add_picture(str(FIG / "impact_estimators.png"), Inches(6.65), Inches(3.3), width=Inches(5.0))
    add_footer(s)

    # 7
    s = blank(prs)
    add_title(s, "Experiment Handoff", "첫 세션 탐색 폭 넛지 A/B 설계")
    add_card(s, 0.75, 1.35, 5.4, 1.25, "Hypothesis", "첫 세션에서 인접 카테고리와 저장 후보를 노출하면 D7 재방문이 오른다.", GREEN, "H1")
    add_card(s, 6.55, 1.35, 5.25, 1.25, "Treatment", "Adjacent categories + save candidates in first session.", TEAL, "Exploration breadth")
    rows = [
        ("Design item", "Spec"),
        ("Randomization", "User-level 50/50 assignment at first session"),
        ("Primary metric", "D7 revisit"),
        ("Guardrails", "buy completion, immediate bounce, session quality, SRM"),
        ("Power", "MDE 2.0pp, 9,805 users/arm, estimated 4.7 weeks"),
    ]
    add_table_like(s, 0.85, 3.1, 11.0, rows, [0.28, 0.72])
    add_footer(s)

    # 8
    s = blank(prs)
    add_title(s, "Decision Rules", "성공/중단/반복 기준을 사전에 둔다")
    rows = [
        ("Decision", "Rule"),
        ("Ship / scale", "D7 revisit lift >= 2.0pp and guardrails not worse"),
        ("Iterate", "Positive but below MDE; inspect segment and UX logs"),
        ("Stop", "No lift or guardrail degradation"),
        ("Invalidate", "SRM or tracking issue"),
        ("Revenue follow-up", "After D7 passes: purchase, repurchase, GMV validation"),
    ]
    add_table_like(s, 0.9, 1.45, 10.9, rows, [0.28, 0.72])
    add_card(s, 0.9, 4.75, 10.9, 1.0, "Why this matters", "실험 전 판정 기준을 고정해야 사후 지표 변경과 해석 흔들림을 줄일 수 있다.", GREEN, "Pre-registered decision logic")
    add_footer(s)

    # 9
    s = blank(prs)
    add_title(s, "What This Demonstrates", "이 프로젝트가 보여주는 데이터 분석가 역량")
    skills = [
        ("1. 문제를 KPI로 바꿨다", "구매완료율 4.4%의 한계를 보고 D7 재방문을 신규 유저 경험의 primary KPI로 재정의"),
        ("2. 데이터 신뢰성을 먼저 잠갔다", "DuckDB QA로 필드, 이벤트 어휘, 기간, 가격, denominator를 검증한 뒤 분석 진행"),
        ("3. 병목과 후보를 분리했다", "조회 -> 찜 최대 유출과 첫날 활동 폭 +27pp 신호를 실험 후보로 연결"),
        ("4. 상관을 인과로 포장하지 않았다", "cart +16.0pp -> +3.2pp, cat>=3 +25.4pp -> -6.5pp로 관측 신호의 한계 명시"),
        ("5. 실행 가능한 실험으로 끝냈다", "D7 revisit, guardrails, MDE 2.0pp, 9,805/arm, 4.7 weeks까지 handoff"),
        ("6. 재현 가능한 산출물로 남겼다", "Makefile pipeline, generated reports, onepager, portfolio report로 검증 경로 제공"),
    ]
    for i, (title, body) in enumerate(skills):
        x = 0.75 + (i % 3) * 4.05
        y = 1.42 + (i // 3) * 2.05
        add_card(s, x, y, 3.55, 1.35, title, body, [GREEN, TEAL, NAVY, AMBER, GREEN, TEAL][i])
    add_card(
        s,
        0.75,
        5.55,
        11.7,
        0.8,
        "Interview takeaway",
        "이 프로젝트의 핵심은 예쁜 차트가 아니라, 낮은 전환율 문제를 재정의하고 관측 신호의 한계를 검증한 뒤 실행 가능한 A/B 설계로 넘긴 사고 과정이다.",
        GREEN,
    )
    add_footer(s)

    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    build()
